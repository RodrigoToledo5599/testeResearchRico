import docx
import os
from pptx import Presentation
from docx import Document
from docx2pdf import convert

# Nada para se ver aqui
# Esse espaço é apenas para eu testar algumas funcionalidades enquanto eu desenvolvo.



class TestingStuff:
    
    def __init__(self, file_name):
        self.file_name = file_name



    def generateNewDocx(self):
        docxpath = os.getcwd()
        datapath = os.path.join(docxpath,f"Data\{self.file_name}")
        returnpath = os.path.join(docxpath,f"Output\{self.file_name}")

        doc = Document(datapath)        
        
        doc.save(returnpath)

    
    
    def generateNewPttx(self):
        
        pptxpath = os.getcwd()
        modelpath = os.path.join(pptxpath,f"ModeloCarteira\{self.file_name}")
        returnpath = os.path.join(pptxpath,f"Output\{self.file_name}")


        pttx = Presentation(modelpath)
        pttx.save(returnpath)


    def getelementFromPPtx(self):
        pptxpath = os.getcwd()
        modelpath = os.path.join(pptxpath,f"ModeloCarteira\{self.file_name}")
        returnpath = os.path.join(pptxpath,f"Output\{self.file_name}")

        pptx = Presentation(modelpath)
        dd = Presentation()
        dd.slides.add_slide(pptx.slides)
        dd.save(returnpath)



    def clearOutputDir(self):
        dd = os.getcwd()
        dir = os.path.join(dd,"Output")
        for file in os.listdir(dir):
            caminho_completo = os.path.join(dir, file)
            print("Excluindo: ",caminho_completo)
            os.remove(caminho_completo)

        
    def word2pdf(self):
        dd = os.getcwd()
        dir = os.path.join(dd,f"Data\Guia de Investimentos (lista de Trends).docx")
        outputDir = os.path.join(os.getcwd(),"/Output")


        convert(dir)
        convert(dir, dd+f"\Data\Guia de Investimentos (lista de Trends).pdf")
        convert(outputDir)


    def gettingDataFromDocx(self):
        docxpath = os.getcwd()
        datapath = os.path.join(docxpath,f"Data\{self.file_name}")
        returnpath = os.path.join(docxpath,f"Output\{self.file_name}")

        i =0
        doc = Document(datapath)
        for parag in doc.paragraphs:
            print(f'{i}: ',parag.text,'\n\n\n-------------------------------------------------------------------------\n\n\n')
            i +=1
            
    def loopingThroughTxt(self):
        lines1 = []
        lines2 = []
        
        with open('WritenContent/textData.txt', 'r') as arquivo:
            for line in arquivo:
                if(line.strip() == "" or line.strip() =="\n"):
                    continue
                if(line.strip() == "Por que está nessa lista de sugestões?"):
                    lines1.append(line)
                if(line.strip() == "" or line.strip() =="\n"):
                    continue
                if(line == "Confira nossa visão para essa classe de ativo aqui.\n"):
                    lines2.append(line)
        print(lines1)
        print(lines2)


if __name__ == "__main__":

    dd = TestingStuff("Guia de Investimentos (lista de Trends).docx")
    dd.loopingThroughTxt()

